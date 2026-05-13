import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    """
    生成AI工单项目分析PPT演示文稿 (中文版)
    输出路径: ../../conclusion/AI工单项目分析.pptx (相对于当前文件位置)
    """
    # 确定输出路径
    current_dir = os.path.dirname(os.path.abspath(__file__))
    output_dir = os.path.join(current_dir, "../../conclusion/_local")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "AI工单项目分析.pptx")

    # 创建演示文稿
    prs = Presentation()

    # --- 幻灯片 1: 封面 ---
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "AI 工单智能分析系统"
    subtitle.text = "项目结构、设计与开发流程分析\n\n生成时间: 2026-02-05"

    # --- 幻灯片 2: 项目概览 ---
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "项目概览"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    text_frame.text = "核心目标: Jira 工单处理的全自动化智能化"
    
    p = text_frame.add_paragraph()
    p.text = "主要价值:"
    p.level = 0
    p.font.bold = True
    
    bullets = [
        "智能分析 (Intelligent Analysis): 从源数据自动提取问题、分类并生成总结。",
        "智能分派 (Smart Dispatch): 基于历史数据推荐处理团队和角色。",
        "由古溯今 (Recall): 召回相似历史问题 (>70% 相似度)，避免重复劳动。",
        "需求规划 (Requirement Planning): 智能生成 PRD、Spec 和原型设计。"
    ]
    for b in bullets:
        p = text_frame.add_paragraph()
        p.text = b
        p.level = 1

    # --- 幻灯片 3: 系统架构 ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "系统技术架构"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    # Backend
    p = text_frame.add_paragraph()
    p.text = "后端服务 (APP/backend):"
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192) # Blue
    
    p = text_frame.add_paragraph()
    p.text = "语言/框架: Python 3 + FastAPI (异步架构)"
    p.level = 1
    p = text_frame.add_paragraph()
    p.text = "核心模块: Analysis (分析), Search (搜索召回), Planning (规划), KB (知识库)"
    p.level = 1
    p = text_frame.add_paragraph()
    p.text = "任务机制: BackgroundTasks 处理长耗时 LLM 任务"
    p.level = 1
    
    # Frontend
    p = text_frame.add_paragraph()
    p.text = "前端应用 (APP/frontend):"
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    
    p = text_frame.add_paragraph()
    p.text = "技术栈: 原生 HTML5 + TailwindCSS + Vanilla JS"
    p.level = 1
    p = text_frame.add_paragraph()
    p.text = "设计理念: 轻量化、免构建、极简部署"
    p.level = 1

    # Storage
    p = text_frame.add_paragraph()
    p.text = "数据持久化:"
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    p = text_frame.add_paragraph()
    p.text = "本地文件系统 (Markdown/JSON/CSV)"
    p.level = 1

    # --- 幻灯片 4: 核心功能模块 ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "核心功能模块"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    items = [
        ("问题智能分析", "批量处理 src/ 目录下的工单，应用业务规则生成结论。"),
        ("周总结分析报告", "自动统计工单趋势、团队负载，生成图表化周报。"),
        ("产品需求智能规划", "Agentic 工作流，基于模板自动生成规格说明书。"),
        ("交互式问答界面", "支持图片理解和语义搜索的类 Chat 界面。")
    ]
    
    for title_text, desc in items:
        p = text_frame.add_paragraph()
        p.text = title_text
        p.font.bold = True
        p.level = 0
        
        p = text_frame.add_paragraph()
        p.text = desc
        p.level = 1

    # --- 幻灯片 5: 开发与设计流程 ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "开发与设计流程规范"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    p = text_frame.add_paragraph()
    p.text = "规则驱动开发 (Rules-Driven):"
    p.font.bold = True
    p = text_frame.add_paragraph()
    p.text = "严格遵守 RULES.md，确保标准化执行。"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "设计优先 (Design-First):"
    p.font.bold = True
    p = text_frame.add_paragraph()
    p.text = "先设计 (design/*.md) 后开发，文档即真理。"
    p.level = 1
    p = text_frame.add_paragraph()
    p.text = "模板驱动生成，保证产物一致性。"
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "Agentic Integration:"
    p.font.bold = True
    p = text_frame.add_paragraph()
    p.text = "Planning -> Execution -> Verification 自修正闭环。"
    p.level = 1

    # --- 幻灯片 6: 项目文件结构 ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "项目文件结构"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    struct_items = [
        "ROOT/",
        "  ├── APP/          # 源代码 (Backend/Frontend, 运行逻辑)",
        "  ├── conclusion/   # 分析结论、报告、wisdom、tests",
        "  ├── design/       # 设计文档、计划、spec模板",
        "  ├── src/          # 原始数据输入 (Jira 导出)",
        "  ├── KB/           # 知识库文件",
        "  └── RULES.md      # 项目核心规则"
    ]
    
    for item in struct_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.name = "Courier New"
        p.font.size = Pt(16)


    # --- 保存 ---
    prs.save(output_path)
    print(f"演示文稿已生成: {output_path}")

if __name__ == "__main__":
    create_presentation()
