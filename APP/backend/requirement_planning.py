"""
需求规划服务 - Requirement Planning Service
用于智能分析需求概要文件并生成规范化的需求文档
"""

import os
import re
import json
import uuid
import threading
from datetime import datetime
from typing import List, Dict, Optional, Any
from dataclasses import dataclass, asdict
from enum import Enum

# Paths - 使用绝对路径
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.normpath(os.path.join(BASE_DIR, "../.."))

SPEC_DIR = os.path.join(PROJECT_ROOT, "design/spec")
TEMPLATE_DIR = os.path.join(PROJECT_ROOT, "design/template")
OUTPUT_DIR = os.path.join(PROJECT_ROOT, "conclusion/requirements")
TEMP_DIR = os.path.join(PROJECT_ROOT, "conclusion/temp")
VERSIONS_DIR = os.path.join(OUTPUT_DIR, "versions")

# Ensure directories exist
for d in [SPEC_DIR, TEMPLATE_DIR, OUTPUT_DIR, TEMP_DIR, VERSIONS_DIR]:
    os.makedirs(d, exist_ok=True)


class TaskStatus(Enum):
    PENDING = "pending"
    RUNNING = "running"
    COMPLETED = "completed"
    FAILED = "failed"
    CANCELLED = "cancelled"


@dataclass
class SpecFileInfo:
    filename: str
    filepath: str
    file_type: str  # "概要需求" or "详细需求" or "其他"
    size_bytes: int
    modified_time: str
    has_output: bool
    output_files: List[str]


@dataclass
class TaskState:
    task_id: str
    spec_file: str
    template: str
    status: TaskStatus
    progress: int  # 0-100
    current_step: str
    created_at: str
    updated_at: str
    output_files: List[str]
    error_message: Optional[str] = None


class RequirementPlanningService:
    def __init__(self, llm_service=None):
        self.llm_service = llm_service
        self.tasks: Dict[str, TaskState] = {}

        # 调试：打印路径
        print(f"[RequirementPlanning] 初始化...")
        print(f"  - PROJECT_ROOT: {PROJECT_ROOT}")
        print(f"  - SPEC_DIR: {SPEC_DIR} (exists: {os.path.exists(SPEC_DIR)})")
        print(f"  - OUTPUT_DIR: {OUTPUT_DIR} (exists: {os.path.exists(OUTPUT_DIR)})")

        self._load_task_state()

    def _load_task_state(self):
        """Load persisted task state from temp directory"""
        state_file = os.path.join(TEMP_DIR, "task_state.json")
        if os.path.exists(state_file):
            try:
                with open(state_file, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    for task_id, task_data in data.items():
                        task_data['status'] = TaskStatus(task_data['status'])
                        self.tasks[task_id] = TaskState(**task_data)
            except Exception as e:
                print(f"Error loading task state: {e}")

    def _save_task_state(self):
        """Persist task state to temp directory"""
        state_file = os.path.join(TEMP_DIR, "task_state.json")
        try:
            data = {}
            for task_id, task in self.tasks.items():
                task_dict = asdict(task)
                task_dict['status'] = task.status.value
                data[task_id] = task_dict
            with open(state_file, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"Error saving task state: {e}")

    def list_spec_files(self) -> List[SpecFileInfo]:
        """List all spec files with their status"""
        files = []
        if not os.path.exists(SPEC_DIR):
            return files

        for filename in os.listdir(SPEC_DIR):
            if not filename.endswith('.md'):
                continue

            filepath = os.path.join(SPEC_DIR, filename)
            stat = os.stat(filepath)

            # Determine file type
            if "概要需求" in filename:
                file_type = "概要需求"
            elif "详细需求" in filename:
                file_type = "详细需求"
            else:
                file_type = "其他"

            # Check for existing outputs
            output_files = self._find_output_files(filename)

            files.append(SpecFileInfo(
                filename=filename,
                filepath=filepath,
                file_type=file_type,
                size_bytes=stat.st_size,
                modified_time=datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M:%S'),
                has_output=len(output_files) > 0,
                output_files=output_files
            ))

        return files

    def _find_output_files(self, spec_filename: str) -> List[str]:
        """Find output files generated from a spec file"""
        outputs = []
        base_name = os.path.splitext(spec_filename)[0]

        if os.path.exists(OUTPUT_DIR):
            for f in os.listdir(OUTPUT_DIR):
                if f.startswith(base_name) and os.path.isfile(os.path.join(OUTPUT_DIR, f)):
                    outputs.append(f)

        return outputs

    def get_file_content(self, filename: str) -> Optional[str]:
        """Read content of a spec file"""
        filepath = os.path.join(SPEC_DIR, filename)
        if not os.path.exists(filepath):
            return None

        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()

    def get_latest_content(self, filename: str) -> dict:
        """
        返回文件的"最新内容"：
        - 若有输出文件，返回按修改时间最新的输出文件内容
        - 若无，返回原始 spec 文件内容
        """
        spec_path = os.path.join(SPEC_DIR, filename)
        if not os.path.exists(spec_path):
            return {"content": None, "source": "spec", "output_filename": None}

        with open(spec_path, 'r', encoding='utf-8') as f:
            original = f.read()

        output_files = self._find_output_files(filename)
        if output_files:
            def mtime(fn):
                return os.path.getmtime(os.path.join(OUTPUT_DIR, fn))
            latest_filename = max(output_files, key=mtime)
            output_path = os.path.join(OUTPUT_DIR, latest_filename)
            try:
                with open(output_path, 'r', encoding='utf-8') as f:
                    latest_content = f.read()
                return {
                    "content": latest_content,
                    "source": "output",
                    "output_filename": latest_filename,
                }
            except Exception:
                pass

        return {
            "content": original,
            "source": "spec",
            "output_filename": None,
        }

    def save_file_content(self, filename: str, content: str) -> bool:
        """Save content to a spec file"""
        filepath = os.path.join(SPEC_DIR, filename)
        try:
            with open(filepath, 'w', encoding='utf-8') as f:
                f.write(content)
            return True
        except Exception as e:
            print(f"Error saving file: {e}")
            return False

    def upload_file(self, filename: str, content: str) -> bool:
        """Upload a new spec file"""
        if not filename.endswith('.md'):
            filename += '.md'
        return self.save_file_content(filename, content)

    def list_templates(self) -> List[Dict[str, str]]:
        """List available templates (md, docx, pptx)"""
        templates = []
        if not os.path.exists(TEMPLATE_DIR):
            return templates

        supported_extensions = ('.md', '.docx', '.pptx')
        
        for filename in os.listdir(TEMPLATE_DIR):
            # Skip temp files (starting with ~)
            if filename.startswith('.') or filename.startswith('~'):
                continue
            
            if not filename.lower().endswith(supported_extensions):
                continue

            # Determine template type
            if "概要" in filename:
                tpl_type = "概要需求"
            elif "详细" in filename:
                tpl_type = "详细需求"
            else:
                tpl_type = "通用"
            
            # Determine format from extension
            ext = os.path.splitext(filename)[1].lower()
            format_map = {'.md': 'Markdown', '.docx': 'Word', '.pptx': 'PPT'}
            file_format = format_map.get(ext, 'Unknown')

            templates.append({
                "filename": filename,
                "type": tpl_type,
                "format": file_format,
                "path": os.path.join(TEMPLATE_DIR, filename)
            })

        return templates

    def get_template_content(self, filename: str) -> Optional[str]:
        """Read content of a template file (supports md, docx, pptx)"""
        filepath = os.path.join(TEMPLATE_DIR, filename)
        if not os.path.exists(filepath):
            return None

        ext = os.path.splitext(filename)[1].lower()
        
        if ext == '.md':
            with open(filepath, 'r', encoding='utf-8') as f:
                return f.read()
        elif ext == '.docx':
            return self._parse_docx_template(filepath)
        elif ext == '.pptx':
            return self._parse_pptx_template(filepath)
        else:
            return None

    def _parse_docx_template(self, filepath: str) -> str:
        """Extract text content from DOCX file as markdown-like format"""
        try:
            from docx import Document
            doc = Document(filepath)
            
            content_parts = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                
                # Convert heading styles to markdown
                style_name = para.style.name if para.style else ''
                if 'Heading 1' in style_name or 'heading 1' in style_name.lower():
                    content_parts.append(f"# {text}")
                elif 'Heading 2' in style_name or 'heading 2' in style_name.lower():
                    content_parts.append(f"## {text}")
                elif 'Heading 3' in style_name or 'heading 3' in style_name.lower():
                    content_parts.append(f"### {text}")
                elif 'Title' in style_name:
                    content_parts.append(f"# {text}")
                else:
                    content_parts.append(text)
            
            # Also extract table content
            for table in doc.tables:
                table_md = self._table_to_markdown(table)
                if table_md:
                    content_parts.append(table_md)
            
            return "\n\n".join(content_parts)
        except Exception as e:
            print(f"Error parsing DOCX: {e}")
            return f"[DOCX解析错误: {e}]"

    def _parse_pptx_template(self, filepath: str) -> str:
        """Extract text content from PPTX file as markdown-like format"""
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            
            content_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = [f"## 幻灯片 {slide_num}"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Check if it's a title
                        if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                            if shape.placeholder_format.type == 1:  # Title
                                slide_content.append(f"### {shape.text.strip()}")
                            else:
                                slide_content.append(shape.text.strip())
                        else:
                            slide_content.append(shape.text.strip())
                    
                    # Handle tables in slides
                    if shape.has_table:
                        table_md = self._pptx_table_to_markdown(shape.table)
                        if table_md:
                            slide_content.append(table_md)
                
                content_parts.append("\n".join(slide_content))
            
            return "\n\n---\n\n".join(content_parts)
        except Exception as e:
            print(f"Error parsing PPTX: {e}")
            return f"[PPTX解析错误: {e}]"

    def _table_to_markdown(self, table) -> str:
        """Convert DOCX table to markdown format"""
        try:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
            
            if len(rows) >= 1:
                # Add header separator after first row
                header_sep = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
                rows.insert(1, header_sep)
            
            return "\n".join(rows)
        except:
            return ""

    def _pptx_table_to_markdown(self, table) -> str:
        """Convert PPTX table to markdown format"""
        try:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
            
            if len(rows) >= 1:
                # Add header separator after first row
                col_count = len(list(table.rows[0].cells))
                header_sep = "| " + " | ".join(["---"] * col_count) + " |"
                rows.insert(1, header_sep)
            
            return "\n".join(rows)
        except:
            return ""

    def start_generation(self,
                         spec_file: str,
                         template: str,
                         output_formats: List[str],
                         final_decision_notes: str,
                         draft_context: Optional[Dict[str, Any]],
                         api_key: str,
                         provider: str = "gemini",
                         model_name: str = "",
                         base_url: str = "") -> str:
        """Start a requirement generation task"""
        task_id = str(uuid.uuid4())[:8]

        task = TaskState(
            task_id=task_id,
            spec_file=spec_file,
            template=template,
            status=TaskStatus.PENDING,
            progress=0,
            current_step="初始化任务",
            created_at=datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            updated_at=datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            output_files=[]
        )

        self.tasks[task_id] = task
        self._save_task_state()

        # Start background thread
        thread = threading.Thread(
            target=self._run_generation,
            args=(task_id, spec_file, template, output_formats, final_decision_notes, draft_context, api_key, provider, model_name, base_url)
        )
        thread.daemon = True
        thread.start()

        return task_id

    def _run_generation(self,
                        task_id: str,
                        spec_file: str,
                        template: str,
                        output_formats: List[str],
                        final_decision_notes: str,
                        draft_context: Optional[Dict[str, Any]],
                        api_key: str,
                        provider: str,
                        model_name: str,
                        base_url: str):
        """Background task for requirement generation"""
        task = self.tasks[task_id]

        try:
            # Step 1: Load spec content
            self._update_task(task_id, TaskStatus.RUNNING, 10, "读取需求概要文件")
            spec_content = self.get_file_content(spec_file)
            if not spec_content:
                raise Exception(f"无法读取文件: {spec_file}")

            # Step 2: Load template
            self._update_task(task_id, TaskStatus.RUNNING, 20, "加载模板文件")
            template_content = self.get_template_content(template)
            if not template_content:
                raise Exception(f"无法读取模板: {template}")

            # Step 3: Call LLM for analysis
            self._update_task(task_id, TaskStatus.RUNNING, 30, "AI智能分析中...")

            if not self.llm_service or not api_key:
                # Fallback: just merge spec with template structure
                generated_content = self._simple_merge(spec_content, template_content, spec_file)
            else:
                generated_content = self._llm_generate(
                    spec_content, template_content, spec_file,
                    final_decision_notes, draft_context, api_key, provider, model_name, base_url
                )

            # Step 4: Run compliance check
            self._update_task(task_id, TaskStatus.RUNNING, 70, "执行合规性检查")
            compliance_result = self._run_compliance_check(generated_content)
            if compliance_result:
                generated_content += f"\n\n---\n\n## 合规性检查结果\n\n{compliance_result}"

            # Step 5: Save outputs
            self._update_task(task_id, TaskStatus.RUNNING, 80, "保存输出文件")
            output_files = self._save_outputs(spec_file, generated_content, output_formats)

            # Step 6: Save version
            self._update_task(task_id, TaskStatus.RUNNING, 90, "保存版本记录")
            self._save_version(spec_file, generated_content)

            # Done
            task.output_files = output_files
            self._update_task(task_id, TaskStatus.COMPLETED, 100, "完成")

        except Exception as e:
            task.error_message = str(e)
            self._update_task(task_id, TaskStatus.FAILED, task.progress, f"失败: {str(e)}")
            # Cleanup any partial output files on failure
            self._cleanup_failed_task(task_id, spec_file)

    def _cleanup_failed_task(self, task_id: str, spec_file: str):
        """Clean up partial files generated by failed task"""
        try:
            # Remove any partial output files for this spec
            base_name = os.path.splitext(spec_file)[0]
            if os.path.exists(OUTPUT_DIR):
                for filename in os.listdir(OUTPUT_DIR):
                    # Match files generated for this spec with recent timestamp
                    if filename.startswith(base_name) and "_生成_" in filename:
                        filepath = os.path.join(OUTPUT_DIR, filename)
                        # Only delete files created in the last 5 minutes (likely from this task)
                        file_mtime = os.path.getmtime(filepath)
                        if datetime.now().timestamp() - file_mtime < 300:
                            os.remove(filepath)
                            print(f"Cleaned up partial file: {filename}")
            
            # Clean old temp files
            if os.path.exists(TEMP_DIR):
                for filename in os.listdir(TEMP_DIR):
                    if filename.startswith(f"task_{task_id}"):
                        filepath = os.path.join(TEMP_DIR, filename)
                        os.remove(filepath)
                        
        except Exception as e:
            print(f"Cleanup error (non-fatal): {e}")

    def _update_task(self, task_id: str, status: TaskStatus, progress: int, step: str):
        """Update task state"""
        if task_id in self.tasks:
            task = self.tasks[task_id]
            task.status = status
            task.progress = progress
            task.current_step = step
            task.updated_at = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            self._save_task_state()

    def _simple_merge(self, spec_content: str, template_content: str, spec_file: str) -> str:
        """Simple merge without LLM - just fill in template placeholders"""
        # Extract requirement name from filename
        req_name = os.path.splitext(spec_file)[0]
        req_name = req_name.replace("概要需求-", "").replace("详细需求-", "")

        result = template_content
        result = result.replace("{requirement_name}", req_name)
        result = result.replace("{author}", "AI生成")
        result = result.replace("{create_date}", datetime.now().strftime('%Y-%m-%d'))
        result = result.replace("{update_date}", datetime.now().strftime('%Y-%m-%d'))

        # Append original spec content
        result += f"\n\n---\n\n## 原始需求概要\n\n{spec_content}"

        return result

    def _llm_generate(self,
                      spec_content: str,
                      template_content: str,
                      spec_file: str,
                      final_decision_notes: str,
                      draft_context: Optional[Dict[str, Any]],
                      api_key: str,
                      provider: str,
                      model_name: str,
                      base_url: str) -> str:
        """Generate requirement document using LLM"""
        req_name = os.path.splitext(spec_file)[0]
        req_name = req_name.replace("概要需求-", "").replace("详细需求-", "")

        # Load ticket data for context (if available)
        ticket_context = self._load_ticket_context(req_name)
        draft_context_block = self._format_draft_context(draft_context)

        prompt = f"""你是一位资深的产品经理，擅长将粗糙的需求概要转化为专业、完整的产品需求文档。

# 任务
请根据下方提供的【需求概要】，按照【文档模板】的结构，生成一份专业、详细、可落地的产品需求文档。

# 需求概要（原始输入，可能不完整）
```
{spec_content[:15000]}
```

# 相关工单数据（用于场景举例和数据推演）
{ticket_context}

# 需求池结构化分析包（优先用于还原业务场景、方案分层、影响面、待确认项）
{draft_context_block}

# 人工最终方案与加工意见
{final_decision_notes or '暂无额外人工拍板意见，请以需求池提交的初稿为主，保守生成终版 PRD。'}

# 文档模板结构
```
{template_content}
```

# 生成要求

## 1. 内容要求
- **深入理解**: 抓住需求的核心思路和目标，即使原始描述粗糙不完整
- **场景补充**: 根据相关工单数据，补充真实的业务场景和用户痛点
- **方案细化**: 将模糊的需求转化为具体的产品方案，包括功能点、交互逻辑
- **数据推演**: 基于工单数据进行数据规模、性能影响的推演分析
- **优先使用结构化分析包**: 如果 `需求池结构化分析包` 中存在 `business_scenarios`、`solution_candidates`、`change_impact`、`pending_questions`、`competitor_comparison`、`analysis_packet`，优先据此生成对应章节，不要忽略这些结构化输入
- **保持模板不变**: 仅填充既有模板章节，不新增或改写模板结构

## 2. 格式要求
- 严格按照模板的章节结构输出，填充每个章节
- 使用 Mermaid 语法生成流程图和架构图，例如：
  ```mermaid
  flowchart TD
      A[开始] --> B{{判断}}
      B -->|是| C[处理]
      B -->|否| D[结束]
  ```
- 使用 Mermaid 生成 ER 图展示数据模型
- 表格使用标准 Markdown 表格语法
- 保持专业的产品文档语言风格

## 3. 质量要求
- 考虑异常流程和边界条件
- 明确数据权限和安全要求
- 提供测试用例设计思路
- 给出风险点和应对策略

# 输出
请直接输出完整的 Markdown 格式需求文档，不要添加额外的说明文字：
"""

        # Collect all chunks from LLM
        result = self.llm_service.call_llm(
            prompt=prompt,
            api_key=api_key,
            provider=provider,
            model_name=model_name,
            base_url=base_url
        )

        # Replace placeholders if any remain in template
        result = result.replace("{requirement_name}", req_name)
        result = result.replace("{author}", "AI智能生成")
        result = result.replace("{create_date}", datetime.now().strftime('%Y-%m-%d'))
        result = result.replace("{update_date}", datetime.now().strftime('%Y-%m-%d'))

        return result

    def _format_draft_context(self, draft_context: Optional[Dict[str, Any]]) -> str:
        """Format requirement-pool draft context for LLM prompt."""
        if not draft_context:
            return "> 当前未提供需求池结构化分析包，请仅基于需求概要与工单上下文生成。"

        summary = {
            "draft_id": draft_context.get("draft_id"),
            "req_id": draft_context.get("req_id"),
            "spec_file": draft_context.get("spec_file"),
            "analysis_summary": draft_context.get("analysis_summary"),
            "core_problem": (draft_context.get("analysis_summary") or {}).get("core_problem"),
            "current_product_behavior": (draft_context.get("analysis_summary") or {}).get("current_product_behavior"),
            "gap_analysis": (draft_context.get("analysis_summary") or {}).get("gap_analysis"),
            "product_layer": (draft_context.get("analysis_summary") or {}).get("product_layer"),
            "product_value": (draft_context.get("analysis_summary") or {}).get("product_value"),
            "ticket_evidence_summary": (draft_context.get("analysis_summary") or {}).get("ticket_evidence_summary", {}),
            "internal_references": (draft_context.get("analysis_summary") or {}).get("internal_references", []),
            "external_references": (draft_context.get("analysis_summary") or {}).get("external_references", []),
            "business_scenarios": draft_context.get("business_scenarios", []),
            "solution_candidates": draft_context.get("solution_candidates", []),
            "change_impact": draft_context.get("change_impact", []),
            "competitor_comparison": draft_context.get("competitor_comparison", []),
            "pending_questions": draft_context.get("pending_questions", []),
            "ticket_summary": draft_context.get("ticket_summary", {}),
            "analysis_packet": draft_context.get("analysis_packet", {}),
        }
        try:
            return json.dumps(summary, ensure_ascii=False, indent=2)
        except Exception as exc:
            return f"> 需求池结构化分析包格式化失败: {exc}"

    def _load_ticket_context(self, req_name: str) -> str:
        """Load relevant ticket data for context"""
        try:
            # Try to find related topic files in conclusion/Topics
            topics_dir = os.path.join(PROJECT_ROOT, "conclusion/_local/Topics")
            context_parts = []
            
            if os.path.exists(topics_dir):
                # Look for files related to the requirement name
                keywords = req_name.split("-") if "-" in req_name else [req_name]
                keywords = [k.strip() for k in keywords if k.strip()]
                
                for filename in os.listdir(topics_dir):
                    if not filename.endswith('.md'):
                        continue
                    # Check if any keyword matches
                    if any(kw in filename for kw in keywords):
                        filepath = os.path.join(topics_dir, filename)
                        with open(filepath, 'r', encoding='utf-8') as f:
                            content = f.read()
                            # Take first 2000 chars as context
                            context_parts.append(f"### 相关工单: {filename}\n{content[:2000]}...")
                        if len(context_parts) >= 3:  # Max 3 related files
                            break
            
            if context_parts:
                return "\n\n".join(context_parts)
            else:
                return "> 未找到相关工单数据，请根据需求概要内容进行分析。"
        except Exception as e:
            return f"> 加载工单数据时出错: {e}"

    def _run_compliance_check(self, content: str) -> str:
        """Run compliance check against best practices"""
        checklist = [
            ("异常流程处理", ["异常", "错误", "失败", "回滚"]),
            ("数据权限定义", ["权限", "角色", "访问控制"]),
            ("性能约束考虑", ["性能", "响应时间", "并发"]),
            ("边界条件覆盖", ["边界", "最大", "最小", "限制"]),
            ("回滚/降级方案", ["回滚", "降级", "备选"]),
            ("测试用例覆盖", ["测试", "用例", "TC"]),
        ]

        results = []
        for check_name, keywords in checklist:
            found = any(kw in content for kw in keywords)
            status = "✅" if found else "⚠️"
            results.append(f"- {status} {check_name}")

        return "\n".join(results)

    def _save_outputs(self, spec_file: str, content: str, formats: List[str]) -> List[str]:
        """Save generated content in requested formats"""
        base_name = os.path.splitext(spec_file)[0]
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_files = []

        for fmt in formats:
            if fmt == "md":
                filename = f"{base_name}_生成_{timestamp}.md"
                filepath = os.path.join(OUTPUT_DIR, filename)
                with open(filepath, 'w', encoding='utf-8') as f:
                    f.write(content)
                output_files.append(filename)

            elif fmt == "docx":
                filename = f"{base_name}_生成_{timestamp}.docx"
                filepath = os.path.join(OUTPUT_DIR, filename)
                self._generate_docx(content, filepath, base_name)
                output_files.append(filename)

            elif fmt == "pptx":
                filename = f"{base_name}_生成_{timestamp}.pptx"
                filepath = os.path.join(OUTPUT_DIR, filename)
                self._generate_pptx(content, filepath, base_name)
                output_files.append(filename)

        return output_files

    def _generate_docx(self, content: str, filepath: str, title: str):
        """Generate DOCX document from markdown content"""
        try:
            from docx import Document
            from docx.shared import Inches, Pt
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            
            doc = Document()
            
            # Add title
            title_para = doc.add_heading(title, level=0)
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Parse markdown content
            lines = content.split('\n')
            current_table = []
            in_code_block = False
            code_content = []
            
            for line in lines:
                # Handle code blocks
                if line.strip().startswith('```'):
                    if in_code_block:
                        # End code block
                        if code_content:
                            para = doc.add_paragraph()
                            run = para.add_run('\n'.join(code_content))
                            run.font.name = 'Courier New'
                            run.font.size = Pt(9)
                        code_content = []
                        in_code_block = False
                    else:
                        in_code_block = True
                    continue
                
                if in_code_block:
                    code_content.append(line)
                    continue
                
                # Handle headers
                if line.startswith('# '):
                    doc.add_heading(line[2:], level=1)
                elif line.startswith('## '):
                    doc.add_heading(line[3:], level=2)
                elif line.startswith('### '):
                    doc.add_heading(line[4:], level=3)
                elif line.startswith('#### '):
                    doc.add_heading(line[5:], level=4)
                # Handle table rows
                elif line.startswith('|'):
                    current_table.append(line)
                else:
                    # Flush table if exists
                    if current_table:
                        self._add_table_to_docx(doc, current_table)
                        current_table = []
                    
                    # Handle list items
                    if line.strip().startswith('- '):
                        doc.add_paragraph(line.strip()[2:], style='List Bullet')
                    elif re.match(r'^\d+\. ', line.strip()):
                        doc.add_paragraph(re.sub(r'^\d+\. ', '', line.strip()), style='List Number')
                    elif line.strip():
                        doc.add_paragraph(line.strip())
            
            # Flush remaining table
            if current_table:
                self._add_table_to_docx(doc, current_table)
            
            doc.save(filepath)
        except Exception as e:
            print(f"Error generating DOCX: {e}")
            # Fallback to text file
            with open(filepath + '.txt', 'w', encoding='utf-8') as f:
                f.write(f"DOCX generation failed: {e}\n\n{content}")

    def _add_table_to_docx(self, doc, table_lines: List[str]):
        """Add a markdown table to DOCX document"""
        try:
            from docx.shared import Pt
            
            # Parse table
            rows = []
            for line in table_lines:
                if '---' in line:
                    continue
                cells = [c.strip() for c in line.split('|')[1:-1]]
                if cells:
                    rows.append(cells)
            
            if not rows:
                return
            
            # Create table
            table = doc.add_table(rows=len(rows), cols=len(rows[0]))
            table.style = 'Table Grid'
            
            for i, row_data in enumerate(rows):
                row = table.rows[i]
                for j, cell_text in enumerate(row_data):
                    if j < len(row.cells):
                        row.cells[j].text = cell_text
                        # Bold header row
                        if i == 0:
                            for para in row.cells[j].paragraphs:
                                for run in para.runs:
                                    run.bold = True
        except Exception as e:
            print(f"Error adding table: {e}")

    def _generate_pptx(self, content: str, filepath: str, title: str):
        """Generate PPTX presentation from markdown content"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = f"AI智能生成\n{datetime.now().strftime('%Y-%m-%d')}"
            
            # Parse content for slides
            sections = self._parse_sections_for_pptx(content)
            
            # Content slide layout
            bullet_slide_layout = prs.slide_layouts[1]
            
            for section_title, section_content in sections:
                slide = prs.slides.add_slide(bullet_slide_layout)
                slide.shapes.title.text = section_title
                
                # Add content
                body_shape = slide.shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()
                
                lines = section_content.strip().split('\n')
                for i, line in enumerate(lines[:8]):  # Max 8 lines per slide
                    line = line.strip()
                    if not line:
                        continue
                    # Clean markdown
                    line = re.sub(r'^[-*]\s*', '', line)
                    line = re.sub(r'^\d+\.\s*', '', line)
                    line = re.sub(r'\*\*([^*]+)\*\*', r'\1', line)
                    line = re.sub(r'\*([^*]+)\*', r'\1', line)
                    
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = line[:100]  # Limit length
                    p.level = 0
            
            prs.save(filepath)
        except Exception as e:
            print(f"Error generating PPTX: {e}")
            with open(filepath + '.txt', 'w', encoding='utf-8') as f:
                f.write(f"PPTX generation failed: {e}\n\n{content}")

    def _parse_sections_for_pptx(self, content: str) -> List[tuple]:
        """Parse markdown content into sections for PPTX slides"""
        sections = []
        current_title = "概述"
        current_content = []
        
        for line in content.split('\n'):
            if line.startswith('## '):
                if current_content:
                    sections.append((current_title, '\n'.join(current_content)))
                current_title = line[3:].strip()
                current_content = []
            elif line.startswith('# '):
                continue  # Skip top-level headers
            else:
                current_content.append(line)
        
        if current_content:
            sections.append((current_title, '\n'.join(current_content)))
        
        return sections[:15]  # Max 15 slides

    def _save_version(self, spec_file: str, content: str):
        """Save version for rollback"""
        base_name = os.path.splitext(spec_file)[0]
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        version_file = f"{base_name}_v{timestamp}.md"
        filepath = os.path.join(VERSIONS_DIR, version_file)

        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(content)

    def get_task_status(self, task_id: str) -> Optional[Dict[str, Any]]:
        """Get task status"""
        if task_id not in self.tasks:
            return None

        task = self.tasks[task_id]
        return {
            "task_id": task.task_id,
            "spec_file": task.spec_file,
            "status": task.status.value,
            "progress": task.progress,
            "current_step": task.current_step,
            "created_at": task.created_at,
            "updated_at": task.updated_at,
            "output_files": task.output_files,
            "error_message": task.error_message
        }

    def cancel_task(self, task_id: str) -> bool:
        """Cancel a running task"""
        if task_id in self.tasks:
            task = self.tasks[task_id]
            if task.status == TaskStatus.RUNNING:
                task.status = TaskStatus.CANCELLED
                self._save_task_state()
                return True
        return False

    def list_versions(self, spec_file: str) -> List[Dict[str, str]]:
        """List all versions of a spec file's outputs"""
        versions = []
        base_name = os.path.splitext(spec_file)[0]

        if os.path.exists(VERSIONS_DIR):
            for f in os.listdir(VERSIONS_DIR):
                if f.startswith(base_name) and f.endswith('.md'):
                    filepath = os.path.join(VERSIONS_DIR, f)
                    stat = os.stat(filepath)
                    versions.append({
                        "filename": f,
                        "created_at": datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M:%S'),
                        "size_bytes": stat.st_size
                    })

        versions.sort(key=lambda x: x['created_at'], reverse=True)

        # 追加原始 spec 文件作为最旧版本（固定在末尾，不参与时间排序）
        spec_path = os.path.join(SPEC_DIR, spec_file)
        if os.path.exists(spec_path):
            stat = os.stat(spec_path)
            versions.append({
                "filename": spec_file,
                "created_at": datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M:%S'),
                "size_bytes": stat.st_size,
                "is_original": True,
                "label": "📌 原始输入文件"
            })

        return versions

    def get_version_content(self, version_file: str) -> Optional[str]:
        """Get content of a specific version"""
        # 先找 versions 目录
        filepath = os.path.join(VERSIONS_DIR, version_file)
        if os.path.exists(filepath):
            with open(filepath, 'r', encoding='utf-8') as f:
                return f.read()
        # fallback：尝试 spec 目录（原始输入文件）
        spec_path = os.path.join(SPEC_DIR, version_file)
        if os.path.exists(spec_path):
            with open(spec_path, 'r', encoding='utf-8') as f:
                return f.read()
        return None

    def refine_content(self, 
                       original_content: str, 
                       user_instruction: str,
                       api_key: str,
                       provider: str = "gemini",
                       model_name: str = "",
                       base_url: str = "") -> str:
        """Refine AI-generated content based on user instructions"""
        if not self.llm_service or not api_key:
            return original_content
        
        prompt = f"""你是一位专业的产品经理助手。请根据用户的修改指令，对以下需求文档进行调整完善。

# 当前文档内容
```markdown
{original_content[:20000]}
```

# 用户修改指令
{user_instruction}

# 要求
1. 仅针对用户指令中提到的部分进行修改
2. 保持文档其他部分不变
3. 保持Markdown格式和Mermaid图表语法
4. 保持专业的产品文档风格

请直接输出修改后的完整文档（Markdown格式）：
"""
        
        result = self.llm_service.call_llm(
            prompt=prompt,
            api_key=api_key,
            provider=provider,
            model_name=model_name,
            base_url=base_url
        )
        
        return result

    def refine_section(self,
                       original_content: str,
                       section_name: str,
                       user_instruction: str,
                       api_key: str,
                       provider: str = "gemini",
                       model_name: str = "",
                       base_url: str = "") -> str:
        """Refine a specific section of the document"""
        if not self.llm_service or not api_key:
            return original_content
        
        prompt = f"""你是一位专业的产品经理助手。请仅针对文档中的"{section_name}"章节，根据用户指令进行修改。

# 当前文档内容
```markdown
{original_content[:20000]}
```

# 目标章节
{section_name}

# 用户修改指令
{user_instruction}

# 要求
1. 仅修改"{section_name}"章节的内容
2. 严格保持其他章节完全不变
3. 保持Markdown格式和Mermaid图表语法
4. 保持专业的产品文档风格

请直接输出修改后的完整文档（Markdown格式）：
"""
        
        result = self.llm_service.call_llm(
            prompt=prompt,
            api_key=api_key,
            provider=provider,
            model_name=model_name,
            base_url=base_url
        )
        
        return result

    def get_output_content(self, filename: str) -> Optional[str]:
        """Get content from output directory"""
        filepath = os.path.join(OUTPUT_DIR, filename)
        if not os.path.exists(filepath):
            return None
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()

    def save_output_content(self, filename: str, content: str) -> bool:
        """Save content to output file and update versions"""
        filepath = os.path.join(OUTPUT_DIR, filename)
        try:
            with open(filepath, 'w', encoding='utf-8') as f:
                f.write(content)
            # Save version
            base_name = os.path.splitext(filename)[0]
            self._save_version(base_name, content)
            return True
        except Exception as e:
            print(f"Error saving output: {e}")
            return False


# Singleton instance
_service_instance = None


def get_requirement_planning_service(llm_service=None):
    global _service_instance
    if _service_instance is None:
        _service_instance = RequirementPlanningService(llm_service)
    return _service_instance
